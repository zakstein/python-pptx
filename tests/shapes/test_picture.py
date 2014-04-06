# encoding: utf-8

"""
Test suite for pptx.shapes.picture module
"""

from __future__ import absolute_import, print_function, unicode_literals

import pytest

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.picture import Picture


class Describe_Picture(object):

    def it_knows_its_shape_type(self, picture):
        assert picture.shape_type == MSO_SHAPE_TYPE.PICTURE

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def picture(self):
        return Picture(None, None)
